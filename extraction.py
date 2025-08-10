"""
PowerPoint content extraction with OCR support and validation.
"""

import logging
import os
import re
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from PIL import Image

# Note: Table class is accessed through shape.table, not imported directly


logger = logging.getLogger(__name__)


class SlideDoc:
    """Represents a single slide with extracted content."""
    
    def __init__(self, slide_num: int, title: str = "", content: str = "", 
                 tables: List[str] = None, image_text: str = ""):
        self.slide_num = slide_num
        self.title = title
        self.content = content
        self.tables = tables or []
        self.image_text = image_text
    
    def get_all_text(self) -> str:
        """Get all text content from the slide."""
        all_text = []
        
        if self.title:
            all_text.append(self.title)
        
        if self.content:
            all_text.append(self.content)
        
        if self.tables:
            all_text.extend(self.tables)
        
        if self.image_text:
            all_text.append(self.image_text)
        
        return "\n".join(all_text)
    
    def __str__(self):
        return f"Slide {self.slide_num}: {self.title[:50]}..."


def validate_inputs(pptx_path: str, img_dir: Optional[str] = None) -> None:
    """Validate input files and directories."""
    # Check PowerPoint file
    pptx_file = Path(pptx_path)
    if not pptx_file.exists():
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
    
    if not pptx_file.suffix.lower() == '.pptx':
        raise ValueError(f"File must be a .pptx file: {pptx_path}")
    
    # Check image directory if provided
    if img_dir:
        img_path = Path(img_dir)
        if not img_path.exists():
            raise FileNotFoundError(f"Image directory not found: {img_dir}")
        
        if not img_path.is_dir():
            raise ValueError(f"Image path must be a directory: {img_dir}")
    
    logger.debug("Input validation passed")


def extract_presentation_content(pptx_path: str, img_dir: Optional[str] = None, 
                               gemini_client=None) -> List[SlideDoc]:
    """Extract comprehensive content from PowerPoint presentation."""
    logger.info("Starting presentation content extraction")
    
    try:
        # Load PowerPoint presentation
        presentation = Presentation(pptx_path)
        slides = []
        
        for i, slide in enumerate(presentation.slides, 1):
            logger.debug(f"Processing slide {i}")
            
            # Extract basic slide content
            slide_doc = extract_slide_content(slide, i)
            
            # Add OCR content from images if available
            if img_dir and gemini_client:
                import asyncio
                image_text = asyncio.run(extract_slide_image_text(i, img_dir, gemini_client))
                if image_text:
                    slide_doc.image_text = image_text
            
            slides.append(slide_doc)
        
        logger.info(f"Extracted content from {len(slides)} slides")
        return slides
    
    except Exception as e:
        logger.error(f"Failed to extract presentation content: {e}")
        raise


def extract_slide_content(slide, slide_num: int) -> SlideDoc:
    """Extract text content from a single slide."""
    title = ""
    content_parts = []
    tables = []
    
    for shape in slide.shapes:
        try:
            # Extract text from text boxes and placeholders
            if hasattr(shape, "text") and shape.text.strip():
                text = clean_text(shape.text)
                
                # Try to identify title vs content
                if not title and (
                    hasattr(shape, 'placeholder_format') and 
                    shape.placeholder_format and 
                    hasattr(shape.placeholder_format, 'type') and
                    shape.placeholder_format.type == 1  # Title placeholder
                ):
                    title = text
                else:
                    content_parts.append(text)
            
            # Extract table content
            elif hasattr(shape, 'has_table') and shape.has_table:
                table_text = extract_table_text(shape.table)
                if table_text:
                    tables.append(table_text)
        
        except Exception as e:
            logger.debug(f"Error processing shape on slide {slide_num}: {e}")
            continue
    
    # If no title was identified, use first content as title
    if not title and content_parts:
        title = content_parts.pop(0)
    
    content = "\n".join(content_parts)
    
    return SlideDoc(
        slide_num=slide_num,
        title=title,
        content=content,
        tables=tables
    )


def extract_table_text(table) -> str:
    """Extract and format text from a PowerPoint table."""
    try:
        table_rows = []
        
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = clean_text(cell.text) if cell.text else ""
                row_cells.append(cell_text)
            
            if any(cell.strip() for cell in row_cells):  # Skip empty rows
                table_rows.append(" | ".join(row_cells))
        
        if table_rows:
            return "\n".join(table_rows)
    
    except Exception as e:
        logger.debug(f"Error extracting table: {e}")
    
    return ""


async def extract_slide_image_text(slide_num: int, img_dir: str, 
                                 gemini_client) -> str:
    """Extract text from slide image using OCR."""
    try:
        # Look for slide image files
        img_path = Path(img_dir)
        possible_names = [
            f"slide{slide_num}.png",
            f"slide{slide_num}.jpg",
            f"slide{slide_num}.jpeg",
            f"Slide{slide_num}.png",
            f"Slide{slide_num}.jpg",
            f"slide_{slide_num}.png",
            f"slide_{slide_num}.jpg"
        ]
        
        image_file = None
        for name in possible_names:
            candidate = img_path / name
            if candidate.exists():
                image_file = candidate
                break
        
        if not image_file:
            logger.debug(f"No image file found for slide {slide_num}")
            return ""
        
        # Extract text using Gemini Vision
        logger.debug(f"Extracting OCR text from {image_file}")
        image_text = await gemini_client.extract_text_from_image(str(image_file))
        
        if image_text:
            logger.debug(f"Extracted {len(image_text)} characters from slide {slide_num} image")
            return clean_text(image_text)
    
    except Exception as e:
        logger.debug(f"Error extracting image text for slide {slide_num}: {e}")
    
    return ""


def clean_text(text: str) -> str:
    """Clean and normalize extracted text."""
    if not text:
        return ""
    
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text)
    
    # Remove common artifacts
    text = re.sub(r'^\d+\s*$', '', text)  # Remove standalone numbers
    text = re.sub(r'^[^\w\s]*$', '', text)  # Remove lines with only punctuation
    
    # Clean up bullet points and formatting
    text = re.sub(r'^[\u2022\u2023\u25E6\u2043\u2219\-\*]+\s*', '• ', text, flags=re.MULTILINE)
    
    return text.strip()


def get_slide_images(img_dir: str) -> List[str]:
    """Get list of available slide image files."""
    if not img_dir:
        return []
    
    img_path = Path(img_dir)
    if not img_path.exists():
        return []
    
    image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp'}
    image_files = []
    
    for file_path in img_path.iterdir():
        if file_path.suffix.lower() in image_extensions:
            image_files.append(str(file_path))
    
    return sorted(image_files)


def extract_numbers_from_text(text: str) -> List[str]:
    """Extract numerical values from text for analysis."""
    # Enhanced number pattern including currency and percentages
    number_pattern = re.compile(
        r'(?:(?:USD?|EUR?|GBP|\$|€|£|Rs\.?)\s*)?'  # Currency prefix
        r'(-?\d{1,3}(?:,\d{3})*(?:\.\d+)?)'         # Number with commas
        r'\s*([KMBTkmbt])?'                          # Suffix (K, M, B, T)
        r'(?:\s*(?:%|percent|USD?|EUR?|GBP|\$|€|£|Rs\.?|hours?|mins?|minutes?|times?|x))?',  # Unit
        re.IGNORECASE
    )
    
    matches = number_pattern.findall(text)
    return [match[0] + (match[1] if match[1] else '') for match in matches if match[0]]


def extract_percentages_from_text(text: str) -> List[float]:
    """Extract percentage values from text."""
    # Pattern for percentages
    percent_pattern = re.compile(r'(\d+(?:\.\d+)?)\s*(?:%|percent)', re.IGNORECASE)
    matches = percent_pattern.findall(text)
    
    percentages = []
    for match in matches:
        try:
            percentages.append(float(match))
        except ValueError:
            continue
    
    return percentages


def extract_dates_from_text(text: str) -> List[str]:
    """Extract date references from text."""
    # Various date patterns
    date_patterns = [
        r'\b\d{1,2}/\d{1,2}/\d{2,4}\b',           # MM/DD/YYYY or DD/MM/YYYY
        r'\b\d{1,2}-\d{1,2}-\d{2,4}\b',           # MM-DD-YYYY or DD-MM-YYYY
        r'\b\d{4}-\d{1,2}-\d{1,2}\b',             # YYYY-MM-DD
        r'\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{1,2},?\s+\d{4}\b',  # Month DD, YYYY
        r'\b\d{1,2}\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4}\b',    # DD Month YYYY
        r'\bQ[1-4]\s+\d{4}\b',                     # Q1 2024
        r'\b\d{4}\b'                              # Just year
    ]
    
    dates = []
    for pattern in date_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        dates.extend(matches)
    
    return dates
